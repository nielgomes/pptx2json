import json
from pptx import Presentation
from tkinter import filedialog
import os

sel_arquivo = filedialog.askopenfilename()
# Abrir o arquivo .ppt
ppt_file = Presentation(sel_arquivo)

# Criar um dicionário vazio para armazenar as formas por slide
shape_dict = {}

# Iterar sobre todas as slides
for i, slide in enumerate(ppt_file.slides):
    # Criar outro dicionário vazio para armazenar as formas e seus textos do slide atual
    shape_text_dict = {}
    # Iterar sobre todas as formas do slide atual
    for shape in slide.shapes:
        # Verificar se a forma tem texto
        if hasattr(shape, 'text'):
            # Extrair o nome da forma e o texto da forma
            shape_name = shape.name
            shape_text = shape.text
            # Dividir o texto da forma em linhas usando o método splitlines()
            shape_lines = shape_text.splitlines()
            # Atribuir ao dicionário interno a chave sendo o nome da forma e o valor sendo uma lista de linhas de texto 
            shape_text_dict[shape_name] = shape_lines 
    # Atribuir ao dicionário externo a chave sendo o número do slide e o valor sendo o dicionário interno 
    shape_dict["Slide"+str("{0:0>2d}".format(i+1))] = shape_text_dict

# Converter o dicionário em uma string JSON usando os parâmetros de formatação 
shape_json = json.dumps(shape_dict, indent=4, sort_keys=True, ensure_ascii=False)

file_path = sel_arquivo
# Extrair o nome do arquivo sem o caminho 
file_name = os.path.basename(file_path)

# Substituir o . por _ 
file_name = file_name.replace('.', '_')

# Juntar o nome modificado com a extensão .json 
json_file = os.path.join(file_name + '.json')

# Escrever a string JSON em um arquivo .json 
with open(json_file, 'w', encoding="utf-8-sig") as f:
    f.write(shape_json)
